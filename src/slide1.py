"""Slide-1 specific rendering.

Expose a `process(prs, folder: Path)` function that the main driver can call.
This module creates a slide that uses `Slide1.jpg` as a background, places two
configurable textboxes, a Lenovo logo at the bottom-right, and a date textbox.
It also still attaches the sales table and any additional images as separate
slides (backward-compatible).
"""
from pathlib import Path
import datetime
import pandas as pd
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image

from .base_slide import read_meta_title, px_to_inches, add_images_from_folder, add_table_slide

# --- Adjustable offsets (in inches) that you can tweak ---
TEXTBOX1_LEFT_IN = 1.0
TEXTBOX1_TOP_IN = 1.0

TEXTBOX2_LEFT_IN = 1.0
TEXTBOX2_TOP_IN = 1.6

DATEBOX_LEFT_IN = 0.5
DATEBOX_TOP_IN = 6.5

# Logo sizing (max width in inches) and margin from edges
LOGO_MAX_WIDTH_IN = 1.2
LOGO_MARGIN_IN = 0.2


def _add_bg_and_boxes(prs, folder: Path):
    """Create the main slide with background image, textboxes, logo, and date."""
    blank = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank)

    # Background image
    bg = folder / 'Slide1.jpg'
    if bg.exists():
        try:
            # scale background to slide width
            prs_width = prs.slide_width
            slide.shapes.add_picture(str(bg), 0, 0, width=prs_width)
        except Exception:
            pass

    # Text Box 1 (Arial, size 24)
    tb1_w = Inches(4)
    tb1_h = Inches(0.8)
    tx1 = slide.shapes.add_textbox(Inches(TEXTBOX1_LEFT_IN), Inches(TEXTBOX1_TOP_IN), tb1_w, tb1_h).text_frame
    tx1.text = "Smarter Technology for All"
    p1 = tx1.paragraphs[0]
    run1 = p1.runs[0]
    run1.font.name = 'Arial'
    run1.font.size = Pt(24)

    # Text Box 2 (Arial Bold, size 72)
    tb2_w = Inches(8)
    tb2_h = Inches(2)
    tx2 = slide.shapes.add_textbox(Inches(TEXTBOX2_LEFT_IN), Inches(TEXTBOX2_TOP_IN), tb2_w, tb2_h).text_frame
    tx2.text = "Lenovo Americas Route Overview"
    p2 = tx2.paragraphs[0]
    run2 = p2.runs[0]
    run2.font.name = 'Arial'
    run2.font.size = Pt(72)
    run2.font.bold = True

    # Date textbox (Arial size 20) with today's date
    today = datetime.date.today()
    date_str = today.strftime('%-d, %B, %Y') if hasattr(today, 'strftime') else today.strftime('%d, %B, %Y')
    # Note: %-d produces no leading zero on some systems; fallback above
    db_w = Inches(4)
    db_h = Inches(0.6)
    dtx = slide.shapes.add_textbox(Inches(DATEBOX_LEFT_IN), Inches(DATEBOX_TOP_IN), db_w, db_h).text_frame
    dtx.text = date_str
    dp = dtx.paragraphs[0]
    drun = dp.runs[0]
    drun.font.name = 'Arial'
    drun.font.size = Pt(20)

    # Lenovo logo at bottom-right
    logo = folder / 'Lenovo-Logo.png'
    if logo.exists():
        try:
            im = Image.open(logo)
            w_px, h_px = im.size
            w_in = px_to_inches(w_px)
            h_in = px_to_inches(h_px)
            scale = min(LOGO_MAX_WIDTH_IN / w_in, 1.0)
            logo_w = Inches(w_in * scale)
            logo_h = Inches(h_in * scale)
            # compute bottom-right position
            left = prs.slide_width - logo_w - Inches(LOGO_MARGIN_IN)
            top = prs.slide_height - logo_h - Inches(LOGO_MARGIN_IN)
            slide.shapes.add_picture(str(logo), left, top, width=logo_w)
        except Exception:
            pass


def _add_sales_table(prs, df: pd.DataFrame, title: str = 'Sales'):
    # Use shared table helper
    add_table_slide(prs, df, title)


def _add_images(prs, folder: Path):
    imgs = list(folder.glob('*.png')) + list(folder.glob('*.jpg')) + list(folder.glob('*.jpeg'))
    imgs = [p for p in imgs if p.name.lower() not in ('slide1.jpg', 'lenovo-logo.png')]
    if not imgs:
        return
    blank = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank)
    add_images_from_folder(slide, folder, exclude_names=['Slide1.jpg', 'Lenovo-Logo.png'])


def process(prs, folder: Path):
    """Render slide 1 from `folder` into `prs`.

    Expected folder contents:
    - meta.txt (optional) - first line is title
    - sales.csv - table to show
    - images (png/jpg) optional
    """
    title = read_meta_title(folder) or folder.name

    # Main designed slide: background, boxes, date, logo
    _add_bg_and_boxes(prs, folder)

    # Add sales table if available
    #sales_csv = folder / 'sales.csv'
    #if sales_csv.exists():
    #    try:
    #        df = pd.read_csv(sales_csv)
    #        _add_sales_table(prs, df, title='Sales Summary')
    #    except Exception:
    #        pass

    # Additional images (not background/logo)
    _add_images(prs, folder)
