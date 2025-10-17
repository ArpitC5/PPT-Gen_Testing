from pathlib import Path
from typing import List, Optional

import pandas as pd
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx import Presentation
from PIL import Image


def read_meta_title(folder: Path) -> Optional[str]:
    meta = folder / 'meta.txt'
    if not meta.exists():
        return None
    try:
        return meta.read_text(encoding='utf-8').strip().splitlines()[0]
    except Exception:
        return None


def px_to_inches(px: int, dpi: int = 96) -> float:
    return px / dpi


def add_table_slide(prs: Presentation, df: pd.DataFrame, title: str = 'Table', *,
                    left=Inches(0.5), top=Inches(0.6), width=Inches(9),
                    title_font_size=Pt(22), table_height=Inches(3.6), insights: Optional[List] = None):
    """Add a slide with a table and optional insights block on the right.

    `insights` is expected to be a list of lines, where each line is a list of
    tuples (text, color) and color is one of 'green','red' or other.
    """
    blank = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank)

    # Title
    tx = slide.shapes.add_textbox(left, Inches(0.2), width, Inches(0.4)).text_frame
    p = tx.paragraphs[0]
    p.text = title
    p.font.size = title_font_size

    rows, cols = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols, left, top + Inches(0.2), width, table_height).table
    # header
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
    # rows
    for i, row in enumerate(df.itertuples(index=False), start=1):
        for j, val in enumerate(row):
            table.cell(i, j).text = str(val)

    # Render insights to the right if provided
    if insights:
        insp_left = Inches(6.0)
        insp_top = Inches(0.6)
        insp_w = Inches(3.0)
        insp_h = Inches(4.0)
        tf = slide.shapes.add_textbox(insp_left, insp_top, insp_w, insp_h).text_frame
        tf.word_wrap = True
        # ensure there's at least one paragraph
        if not tf.paragraphs:
            tf.add_paragraph()
        for line in insights:
            p = tf.add_paragraph()
            for piece in line:
                text, color = piece
                run = p.add_run()
                run.text = text
                font = run.font
                font.size = Pt(12)
                if color == 'green':
                    font.color.rgb = RGBColor(0x00, 0x80, 0x00)
                elif color == 'red':
                    font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
                else:
                    font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    return slide


def add_images_from_folder(slide, folder: Path, *, max_width=Inches(6), left=Inches(1), top=Inches(1), exclude_names: Optional[List[str]] = None):
    imgs = sorted(folder.glob('*.png')) + sorted(folder.glob('*.jpg')) + sorted(folder.glob('*.jpeg'))
    if exclude_names:
        lower_ex = {n.lower() for n in exclude_names}
        imgs = [p for p in imgs if p.name.lower() not in lower_ex]
    y = top
    for img in imgs:
        try:
            im = Image.open(img)
            w_px, h_px = im.size
            w_in = px_to_inches(w_px)
            h_in = px_to_inches(h_px)
            scale = min(max_width.inches / w_in, 1.0) if w_in > 0 else 1.0
            slide.shapes.add_picture(str(img), left, y, width=Inches(w_in * scale))
            y = y + Inches(h_in * scale) + Inches(0.2)
        except Exception:
            continue
