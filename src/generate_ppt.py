#Usage: python -m src.generate_ppt <slides_root> <output.pptx>

import sys
import os
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image


def add_title_slide(prs, title_text, subtitle_text=None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    if subtitle_text:
        subtitle.text = subtitle_text


def add_table_slide(prs, title, df: pd.DataFrame):
    # Use a blank layout
    blank = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank)
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    # title
    tx = slide.shapes.add_textbox(left, Inches(0.2), width, Inches(0.5)).text_frame
    p = tx.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    rows, cols = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols, left, top + Inches(0.6), width, Inches(4)).table

    # header
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True

    for i, row in enumerate(df.itertuples(index=False), start=1):
        for j, val in enumerate(row):
            table.cell(i, j).text = str(val)


def add_images_from_folder(slide, folder: Path, max_width=Inches(4), left=Inches(5.5), top=Inches(1.5)):
    imgs = list(folder.glob('*.png')) + list(folder.glob('*.jpg')) + list(folder.glob('*.jpeg'))
    y = top
    for img in imgs:
        try:
            im = Image.open(img)
            width, height = im.size
            # convert px to inches approx assuming 96 dpi
            px_to_in = lambda px: px / 96
            w_in = px_to_in(width)
            h_in = px_to_in(height)
            scale = min(max_width.inches / w_in, 1.0)
            pic = slide.shapes.add_picture(str(img), left, y, width=Inches(w_in * scale))
            y = y + Inches(h_in * scale) + Inches(0.2)
        except Exception:
            continue


def process_slide_folder(prs, folder: Path):
    # Expect a 'meta.txt' for title (optional), csv files for tables, images
    title = folder.name
    meta = folder / 'meta.txt'
    if meta.exists():
        try:
            title = meta.read_text(encoding='utf-8').strip().splitlines()[0]
        except Exception:
            pass

    # Add a simple title slide for the first slide
    add_title_slide(prs, title)

    # Add any CSV as tables
    for csv in folder.glob('*.csv'):
        try:
            df = pd.read_csv(csv)
            add_table_slide(prs, csv.stem, df)
        except Exception:
            continue

    # Add images
    blank = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank)
    add_images_from_folder(slide, folder)


"""Usage: python -m src.generate_ppt <slides_root> <output.pptx>

This module is the generic driver. Slide-specific rendering is delegated to
per-slide modules (for example `src.slide1`).
"""

import sys
from pathlib import Path
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image

# Import slide-specific handlers
try:
    from . import slide1
except Exception:
    slide1 = None
try:
    from . import slide2
except Exception:
    slide2 = None


def add_title_slide(prs, title_text, subtitle_text=None):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    try:
        subtitle = slide.placeholders[1]
    except Exception:
        subtitle = None
    title.text = title_text
    if subtitle_text and subtitle is not None:
        subtitle.text = subtitle_text


def process_generic_folder(prs, folder: Path):
    # Generic fallback processor: title, CSV tables, and an image slide.
    title = folder.name
    meta = folder / 'meta.txt'
    if meta.exists():
        try:
            title = meta.read_text(encoding='utf-8').strip().splitlines()[0]
        except Exception:
            pass

    add_title_slide(prs, title)

    # Add any CSV as tables (simple rendering)
    for csv in folder.glob('*.csv'):
        try:
            df = pd.read_csv(csv)
            # minimal table: reuse a simple layout
            blank = prs.slide_layouts[5]
            slide = prs.slides.add_slide(blank)
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            tx = slide.shapes.add_textbox(left, Inches(0.2), width, Inches(0.5)).text_frame
            p = tx.paragraphs[0]
            p.text = csv.stem
            p.font.size = Pt(20)

            rows, cols = df.shape[0] + 1, df.shape[1]
            table = slide.shapes.add_table(rows, cols, left, top + Inches(0.6), width, Inches(4)).table
            for j, col in enumerate(df.columns):
                table.cell(0, j).text = str(col)
            for i, row in enumerate(df.itertuples(index=False), start=1):
                for j, val in enumerate(row):
                    table.cell(i, j).text = str(val)
        except Exception:
            continue

    # Add images to a separate blank slide
    imgs = list(folder.glob('*.png')) + list(folder.glob('*.jpg')) + list(folder.glob('*.jpeg'))
    if imgs:
        blank = prs.slide_layouts[5]
        slide = prs.slides.add_slide(blank)
        y = Inches(1)
        for img in imgs:
            try:
                im = Image.open(img)
                w, h = im.size
                px_to_in = lambda px: px / 96
                w_in = px_to_in(w)
                h_in = px_to_in(h)
                max_w = Inches(6)
                scale = min(max_w.inches / w_in, 1.0)
                slide.shapes.add_picture(str(img), Inches(1), y, width=Inches(w_in * scale))
                y = y + Inches(h_in * scale) + Inches(0.2)
            except Exception:
                continue


def generate(slides_root: str, output: str):
    prs = Presentation()
    slides_root = Path(slides_root)
    if not slides_root.exists():
        raise SystemExit(f"Slides root {slides_root} does not exist")

    # Process folders in sorted order
    for folder in sorted([p for p in slides_root.iterdir() if p.is_dir()]):
        # slide2 takes precedence if available
        if folder.name == 'slide2' and slide2 is not None:
            try:
                slide2.process(prs, folder)
                continue
            except Exception:
                # fallback to generic if slide2 fails
                pass

        if folder.name == 'slide1' and slide1 is not None:
            try:
                slide1.process(prs, folder)
                continue
            except Exception:
                # fallback to generic if slide1 fails
                pass

        process_generic_folder(prs, folder)

    prs.save(output)
    print(f"Saved presentation to {output}")


def main(argv):
    if len(argv) < 3:
        print(__doc__)
        raise SystemExit("Usage: python -m src.generate_ppt <slides_root> <output.pptx>")
    generate(argv[1], argv[2])


if __name__ == '__main__':
    main(sys.argv)
