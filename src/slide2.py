"""Slide-2 renderer: loads Target-V-Actual.xlsx and generates a table plus insights.

Expected files in `slides/slide2/`:
- Target-V-Actual.xlsx - workbook with a table containing columns:
  Country, Target, Actual, % to Asp, YoY
"""
from pathlib import Path
import re
import pandas as pd
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ZERO_WIDTH = '\u200b'


def _clean_header(h: str) -> str:
    if not isinstance(h, str):
        return h
    return h.replace(ZERO_WIDTH, '').strip()


def _parse_currency(s):
    if pd.isna(s):
        return 0.0
    if isinstance(s, (int, float)):
        return float(s)
    # remove any non-digit except dot and minus
    cleaned = re.sub(r"[^0-9\.-]", "", str(s))
    try:
        return float(cleaned) if cleaned not in ('', '.', '-') else 0.0
    except Exception:
        return 0.0


def _parse_percent(s):
    if pd.isna(s):
        return 0.0
    if isinstance(s, (int, float)):
        return float(s)
    txt = str(s).replace(ZERO_WIDTH, '').strip()
    txt = txt.replace('%', '')
    try:
        return float(txt)
    except Exception:
        return 0.0


def _add_table_slide(prs, df: pd.DataFrame, title: str = 'Target vs Actual', insights=None):
    blank = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank)
    left = Inches(0.5)
    top = Inches(0.6)
    width = Inches(9)

    # Title
    tx = slide.shapes.add_textbox(left, Inches(0.2), width, Inches(0.4)).text_frame
    p = tx.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)

    rows, cols = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols, left, top + Inches(0.2), width, Inches(3.6)).table
    # header
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
    # rows
    for i, row in enumerate(df.itertuples(index=False), start=1):
        for j, val in enumerate(row):
            table.cell(i, j).text = str(val)

    # If insights provided, render them on the same slide (right side)
    if insights:
        insp_left = Inches(6.0)
        insp_top = Inches(0.6)
        insp_w = Inches(3.0)
        insp_h = Inches(4.0)
        tf = slide.shapes.add_textbox(insp_left, insp_top, insp_w, insp_h).text_frame
        tf.word_wrap = True
        if not tf.paragraphs:
            tf.add_paragraph()
        for line in insights:
            p = tf.add_paragraph()
            for piece in line:
                text, color = piece
                run = p.add_run()
                run.text = text
                font = run.font
                font.size = Pt(12)
                if color == 'green':
                    font.color.rgb = RGBColor(0x00, 0x80, 0x00)
                elif color == 'red':
                    font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
                else:
                    font.color.rgb = RGBColor(0x00, 0x00, 0x00)


def _format_currency(val):
    try:
        return f"{val:,.0f}"
    except Exception:
        return str(val)


def _add_insights_slide(prs, insights):
    blank = prs.slide_layouts[5]
    slide = prs.slides.add_slide(blank)
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(4)
    tf = slide.shapes.add_textbox(left, top, width, height).text_frame
    tf.word_wrap = True
    # Ensure initial paragraph exists
    if not tf.paragraphs:
        tf.add_paragraph()
    for line in insights:
        p = tf.add_paragraph()
        for piece in line:
            text, color = piece
            run = p.add_run()
            run.text = text
            font = run.font
            font.size = Pt(14)
            if color == 'green':
                font.color.rgb = RGBColor(0x00, 0x80, 0x00)
            elif color == 'red':
                font.color.rgb = RGBColor(0xC0, 0x00, 0x00)
            else:
                font.color.rgb = RGBColor(0x00, 0x00, 0x00)


def process(prs, folder: Path):
    title = folder.name
    meta = folder / 'meta.txt'
    if meta.exists():
        try:
            title = meta.read_text(encoding='utf-8').strip().splitlines()[0]
        except Exception:
            pass

    excel = folder / 'Target-V-Actual.xlsx'
    if not excel.exists():
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        try:
            slide.shapes.title.text = title
        except Exception:
            pass
        return

    try:
        df = pd.read_excel(excel)
    except Exception:
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        try:
            slide.shapes.title.text = title + ' (failed to read Excel)'
        except Exception:
            pass
        return

    # Clean headers containing invisible characters
    df.columns = [_clean_header(c) for c in df.columns]

    # Standardize column names we expect
    col_map = {}
    for c in df.columns:
        low = c.lower()
        if 'country' in low:
            col_map[c] = 'Country'
        elif 'target' in low:
            col_map[c] = 'Target'
        elif 'actual' in low:
            col_map[c] = 'Actual'
        elif 'to asp' in low or '% to asp' in low or 'asp' in low:
            col_map[c] = '% to Asp'
        elif 'yoy' in low or 'y-o-y' in low:
            col_map[c] = 'YoY'

    df = df.rename(columns=col_map)

    # Parse numeric columns robustly
    if 'Target' in df.columns:
        df['Target'] = df['Target'].apply(_parse_currency)
    if 'Actual' in df.columns:
        df['Actual'] = df['Actual'].apply(_parse_currency)
    if '% to Asp' in df.columns:
        df['% to Asp'] = df['% to Asp'].apply(_parse_percent)
    if 'YoY' in df.columns:
        df['YoY'] = df['YoY'].apply(_parse_percent)

    # Prepare display values
    display_df = df.copy()
    if 'Target' in display_df.columns:
        display_df['Target'] = display_df['Target'].apply(_format_currency)
    if 'Actual' in display_df.columns:
        display_df['Actual'] = display_df['Actual'].apply(_format_currency)
    if '% to Asp' in display_df.columns:
        display_df['% to Asp'] = display_df['% to Asp'].apply(lambda v: f"{v:.1f}%")
    if 'YoY' in display_df.columns:
        display_df['YoY'] = display_df['YoY'].apply(lambda v: f"{v:+.1f}%")

    # create a single combined slide with table and insights (done below)

    # Generate insights
    insights = []
    if 'Target' in df.columns and 'Actual' in df.columns:
        total_target = df['Target'].sum()
        total_actual = df['Actual'].sum()
        pct = (total_actual / total_target * 100) if total_target else 0
        diff = total_actual - total_target
        insights.append([
            (f"Total Target: ", 'black'),
            (f"{_format_currency(total_target)}  ", 'black'),
            (f"Total Actual: ", 'black'),
            (f"{_format_currency(total_actual)}  ", 'green' if diff>=0 else 'red'),
            (f"({pct:.1f}% of target)\n", 'black'),
        ])

    if '% to Asp' in df.columns and not df['% to Asp'].isnull().all():
        sorted_by_pct = df.sort_values('% to Asp', ascending=False)
        top = sorted_by_pct.iloc[0]
        bottom = sorted_by_pct.iloc[-1]
        insights.append([
            (f"Top % to Asp: {top['Country']} {top['% to Asp']:.1f}%\n", 'green' if top['% to Asp']>=0 else 'red')
        ])
        insights.append([
            (f"Lowest % to Asp: {bottom['Country']} {bottom['% to Asp']:.1f}%\n", 'red' if bottom['% to Asp']<0 else 'black')
        ])

    if 'YoY' in df.columns and not df['YoY'].isnull().all():
        sorted_yoy = df.sort_values('YoY', ascending=False)
        best = sorted_yoy.iloc[0]
        worst = sorted_yoy.iloc[-1]
        insights.append([
            (f"Top YoY: {best['Country']} {best['YoY']:+.1f}%\n", 'green' if best['YoY']>=0 else 'red')
        ])
        insights.append([
            (f"Worst YoY: {worst['Country']} {worst['YoY']:+.1f}%\n", 'red' if worst['YoY']<0 else 'black')
        ])

    # Add insights to the same table slide by re-creating the table slide with insights
    _add_table_slide(prs, display_df, title='Target vs Actual', insights=insights)
